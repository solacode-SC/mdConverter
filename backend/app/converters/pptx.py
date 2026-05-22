from pptx import Presentation

def convert_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    markdown_lines = []
    
    for i, slide in enumerate(prs.slides):
        markdown_lines.append(f"## Slide {i + 1}")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # simplistic heuristic for title vs body
                    if shape == slide.shapes[0]:
                        markdown_lines.append(f"### {text}")
                    else:
                        markdown_lines.append(f"- {text}")
        markdown_lines.append("")
        
    return "\n".join(markdown_lines)
